import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches


def create_chart(total_df, executed_df, goat_df, tdm_df, node, filename):

    t1 = total_df[total_df['node_id'] == node]
    t2 = executed_df[executed_df['node_id'] == node]
    t3 = goat_df[goat_df['node_id'] == node]
    t4 = tdm_df[tdm_df['node_id'] == node]

    merged = t1.merge(t2, on=['date','node_id'], how='left') \
               .merge(t3, on=['date','node_id'], how='left') \
               .merge(t4, on=['date','node_id'], how='left')

    merged = merged.fillna(0)

    x = np.arange(len(merged['date']))
    width = 0.2

    plt.figure(figsize=(10,5))

    bars1 = plt.bar(x - 1.5*width, merged['total_count'], width, label="Total", color='blue')
    bars2 = plt.bar(x - 0.5*width, merged['executed_count'], width, label="Executed", color='green')
    bars3 = plt.bar(x + 0.5*width, merged['goat_count'], width, label="GOAT Finished", color='orange')
    bars4 = plt.bar(x + 1.5*width, merged['tdm_count'], width, label="TDM Finished", color='red')

    plt.xticks(x, merged['date'].astype(str), rotation=45)
    plt.xlabel("Date")
    plt.ylabel("Count")
    plt.title(f"{node} Day-wise Statistics")
    plt.legend()
    plt.grid(axis='y')

    for bars in [bars1, bars2, bars3, bars4]:
        for bar in bars:
            height = bar.get_height()
            plt.text(bar.get_x()+bar.get_width()/2,
                     height,
                     int(height),
                     ha='center',
                     va='bottom')

    plt.tight_layout()
    plt.savefig(filename)
    plt.close()


def main():

    df = pd.read_csv("Data.csv")

    df['since'] = pd.to_datetime(df['since'], errors='coerce')
    df['date'] = df['since'].dt.date

    # -----------------------------
    # TOTAL COUNT
    # -----------------------------
    total_df = (
        df.groupby(['date','node_id'])['group_id']
        .count()
        .reset_index()
        .rename(columns={'group_id':'total_count'})
    )

    # -----------------------------
    # EXECUTED
    # -----------------------------
    executed_df = (
        df[df['tasker_status']=="EXECUTED"]
        .groupby(['date','node_id'])['group_id']
        .count()
        .reset_index()
        .rename(columns={'group_id':'executed_count'})
    )

    # -----------------------------
    # GOAT FINISHED
    # -----------------------------
    goat_df = (
        df[df['goat_status']=="FINISHED"]
        .groupby(['date','node_id'])['group_id']
        .count()
        .reset_index()
        .rename(columns={'group_id':'goat_count'})
    )

    # -----------------------------
    # TDM FINISHED
    # -----------------------------
    tdm_df = (
        df[df['tdm_gen_status']=="FINISHED"]
        .groupby(['date','node_id'])['group_id']
        .count()
        .reset_index()
        .rename(columns={'group_id':'tdm_count'})
    )

    nodes = df['node_id'].unique()

    prs = Presentation()

    for node in nodes:

        image_name = f"{node}_combined.png"

        create_chart(
            total_df,
            executed_df,
            goat_df,
            tdm_df,
            node,
            image_name
        )

        # -----------------------------
        # CREATE SLIDE
        # -----------------------------
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{node} Day-wise Statistics"

        slide.shapes.add_picture(
            image_name,
            Inches(1),
            Inches(1.5),
            width=Inches(8)
        )

    prs.save("Node_Report.pptx")

    print("PPT created successfully!")


if __name__ == "__main__":
    main()