# # import pandas as pd
# # import matplotlib.pyplot as plt
# # from pptx import Presentation
# # from pptx.util import Inches

# # def main():
# #     # Load CSV
# #     df = pd.read_csv("Data.csv")

# #     # Convert 'since' to datetime
# #     df['since'] = pd.to_datetime(df['since'], dayfirst=True, errors='coerce')

# #     # Extract date only
# #     df['date'] = df['since'].dt.date

# #     # Group by date & node, count group_id
# #     result = (
# #         df.groupby(['date', 'node_id'])['group_id']
# #         .count()
# #         .reset_index()
# #         .rename(columns={'group_id': 'count_group_id'})
# #     )

# #     # Save to Excel
# #     result.to_excel("output.xlsx", index=False)
# #     print(" Result saved to output.xlsx")

# #     # -------- BAR GRAPH WITH COUNT LABELS --------
# #     nodes = result['node_id'].unique()
# #     image_files = []   # to store chart paths

# #     for node in nodes:
# #         node_df = result[result['node_id'] == node]

# #         plt.figure(figsize=(10, 5))
# #         bars = plt.bar(node_df['date'].astype(str), node_df['count_group_id'])

# #         plt.title(f"Date-wise Count of group_id for {node}")
# #         plt.xlabel("Date")
# #         plt.ylabel("Count")
# #         plt.grid(axis='y')
# #         plt.xticks(rotation=45)

# #         # Add labels on each bar
# #         for bar in bars:
# #             height = bar.get_height()
# #             plt.text(
# #                 bar.get_x() + bar.get_width() / 2,
# #                 height,
# #                 str(height),
# #                 ha='center',
# #                 va='bottom',
# #                 fontsize=9,
# #                 fontweight='bold'
# #             )

# #         plt.tight_layout()
# #         filename = f"{node}_bar_chart.png"
# #         plt.savefig(filename)
# #         plt.close()

# #         image_files.append((node, filename))
# #         print(f" Saved bar graph: {filename}")

# #     # -------- CREATE PPT WITH BOTH IMAGES ON ONE SLIDE --------
# #     prs = Presentation()
# #     slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content layout

# #     # Title
# #     slide.shapes.title.text = "Group Count Trends for Nodes"

# #     # Left image
# #     left_img = image_files[0][1]   # ElSauce1
# #     slide.shapes.add_picture(left_img, Inches(0.3), Inches(1.5), width=Inches(4.8))

# #     # Right image
# #     right_img = image_files[1][1]  # ElSauce2
# #     slide.shapes.add_picture(right_img, Inches(5.1), Inches(1.5), width=Inches(4.8))

# #     prs.save("output.pptx")
# #     print(" PPT created: output.pptx with both images on FIRST PAGE")

# #     print(" All tasks completed successfully!")

# # if __name__ == "__main__":
# #     main()



# import pandas as pd
# import matplotlib.pyplot as plt
# from pptx import Presentation
# from pptx.util import Inches

# def main():

#     # Load CSV
#     df = pd.read_csv("Data.csv")

#     # Convert datetime
#     df['since'] = pd.to_datetime(df['since'], dayfirst=True, errors='coerce')
#     df['date'] = df['since'].dt.date

#     # ---------------------------------------------
#     # 1️⃣ GROUP_ID COUNT (DATE + NODE)
#     # ---------------------------------------------
#     group_result = (
#         df.groupby(['date', 'node_id'])['group_id']
#         .count()
#         .reset_index()
#         .rename(columns={'group_id': 'count_group_id'})
#     )

#     group_result.to_excel("group_count.xlsx", index=False)

#     nodes = group_result['node_id'].unique()
#     group_images = []

#     for node in nodes:
#         node_df = group_result[group_result['node_id'] == node]

#         plt.figure(figsize=(10,5))
#         bars = plt.bar(node_df['date'].astype(str), node_df['count_group_id'])

#         plt.title(f"Total Tasks per Day - {node}")
#         plt.xlabel("Date")
#         plt.ylabel("Total Tasks")
#         plt.xticks(rotation=45)
#         plt.grid(axis='y')

#         for bar in bars:
#             height = bar.get_height()
#             plt.text(bar.get_x()+bar.get_width()/2,
#                      height,
#                      str(height),
#                      ha='center',
#                      va='bottom',
#                      fontweight='bold')

#         plt.tight_layout()

#         filename = f"{node}_group_count.png"
#         plt.savefig(filename)
#         plt.close()

#         group_images.append(filename)

#     # ---------------------------------------------
#     # 2️⃣ EXECUTED TASKER STATUS COUNT
#     # ---------------------------------------------
#     executed_df = df[df['tasker_status'] == "EXECUTED"]

#     executed_result = (
#         executed_df.groupby(['date', 'node_id'])['group_id']
#         .count()
#         .reset_index()
#         .rename(columns={'group_id': 'executed_count'})
#     )

#     executed_result.to_excel("executed_count.xlsx", index=False)

#     executed_images = []

#     for node in nodes:
#         node_df = executed_result[executed_result['node_id'] == node]

#         plt.figure(figsize=(10,5))
#         bars = plt.bar(node_df['date'].astype(str), node_df['executed_count'])

#         plt.title(f"Executed Tasks per Day - {node}")
#         plt.xlabel("Date")
#         plt.ylabel("Executed Tasks")
#         plt.xticks(rotation=45)
#         plt.grid(axis='y')

#         for bar in bars:
#             height = bar.get_height()
#             plt.text(bar.get_x()+bar.get_width()/2,
#                      height,
#                      str(height),
#                      ha='center',
#                      va='bottom',
#                      fontweight='bold')

#         plt.tight_layout()

#         filename = f"{node}_executed_count.png"
#         plt.savefig(filename)
#         plt.close()

#         executed_images.append(filename)

#     # ---------------------------------------------
#     # 3️⃣ CREATE PPT
#     # ---------------------------------------------
#     prs = Presentation()

#     # -----------------------------
#     # PAGE 1 → TOTAL TASK COUNT
#     # -----------------------------
#     slide1 = prs.slides.add_slide(prs.slide_layouts[5])
#     slide1.shapes.title.text = "Total Task Count (Date vs Node)"

#     slide1.shapes.add_picture(group_images[0], Inches(0.3), Inches(1.5), width=Inches(4.8))
#     slide1.shapes.add_picture(group_images[1], Inches(5.1), Inches(1.5), width=Inches(4.8))


#     # -----------------------------
#     # PAGE 2 → EXECUTED TASK COUNT
#     # -----------------------------
#     slide2 = prs.slides.add_slide(prs.slide_layouts[5])
#     slide2.shapes.title.text = "Executed Task Count (Date vs Node)"

#     slide2.shapes.add_picture(executed_images[0], Inches(0.3), Inches(1.5), width=Inches(4.8))
#     slide2.shapes.add_picture(executed_images[1], Inches(5.1), Inches(1.5), width=Inches(4.8))


#     prs.save("output.pptx")

#     print("Excel files created")
#     print("Graphs generated")
#     print("PPT created with 2 pages")

# if __name__ == "__main__":
#     main()



import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

def create_bar_chart(df, value_column, title, ylabel, filename_prefix):

    nodes = df['node_id'].unique()
    image_files = []

    for node in nodes:
        node_df = df[df['node_id'] == node]

        plt.figure(figsize=(10,5))
        bars = plt.bar(node_df['date'].astype(str), node_df[value_column])

        plt.title(f"{title} - {node}")
        plt.xlabel("Date")
        plt.ylabel(ylabel)
        plt.grid(axis='y')
        plt.xticks(rotation=45)

        for bar in bars:
            height = bar.get_height()
            plt.text(bar.get_x()+bar.get_width()/2,
                     height,
                     str(height),
                     ha='center',
                     va='bottom',
                     fontweight='bold')

        plt.tight_layout()

        filename = f"{node}_{filename_prefix}.png"
        plt.savefig(filename)
        plt.close()

        image_files.append(filename)

    return image_files


def main():

    df = pd.read_csv("Data.csv")

    df['since'] = pd.to_datetime(df['since'], dayfirst=True, errors='coerce')
    df['date'] = df['since'].dt.date

    # ------------------------------------------------
    # 1 TOTAL TASK COUNT
    # ------------------------------------------------
    total_df = (
        df.groupby(['date','node_id'])['group_id']
        .count()
        .reset_index()
        .rename(columns={'group_id':'total_count'})
    )

    total_images = create_bar_chart(
        total_df,
        "total_count",
        "Total Tasks",
        "Total Tasks",
        "total_tasks"
    )

    # ------------------------------------------------
    # 2 TASKER_STATUS EXECUTED
    # ------------------------------------------------
    executed_df = df[df['tasker_status'] == "EXECUTED"]

    executed_df = (
        executed_df.groupby(['date','node_id'])['group_id']
        .count()
        .reset_index()
        .rename(columns={'group_id':'executed_count'})
    )

    executed_images = create_bar_chart(
        executed_df,
        "executed_count",
        "Executed Tasks",
        "Executed Count",
        "executed_tasks"
    )

    # ------------------------------------------------
    # 3 GOAT_STATUS FINISHED
    # ------------------------------------------------
    finished_df = df[df['goat_status'] == "FINISHED"]

    finished_df = (
        finished_df.groupby(['date','node_id'])['group_id']
        .count()
        .reset_index()
        .rename(columns={'group_id':'finished_count'})
    )

    finished_images = create_bar_chart(
        finished_df,
        "finished_count",
        "GOAT Finished Tasks",
        "Finished Count",
        "goat_finished"
    )

    # ------------------------------------------------
    # 4 TDM_GEN_STATUS FINISHED
    # ------------------------------------------------
    tdm_finished_df = df[df['tdm_gen_status'] == "FINISHED"]

    tdm_finished_df = (
    tdm_finished_df.groupby(['date','node_id'])['group_id']
    .count()
    .reset_index()
    .rename(columns={'group_id':'tdm_finished_count'})
    )

    tdm_finished_images = create_bar_chart(
    tdm_finished_df,
    "tdm_finished_count",
    "TDM Generated FINISHED Tasks",
    "Finished Count",
    "tdm_finished"
    )

    # ------------------------------------------------
    # CREATE PPT
    # ------------------------------------------------
    prs = Presentation()

    # PAGE 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    slide1.shapes.title.text = "Total Task Count"
    slide1.shapes.add_picture(total_images[0], Inches(0.3), Inches(1.5), width=Inches(4.8))
    slide1.shapes.add_picture(total_images[1], Inches(5.1), Inches(1.5), width=Inches(4.8))

    # PAGE 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Tasker Status EXECUTED Count"
    slide2.shapes.add_picture(executed_images[0], Inches(0.3), Inches(1.5), width=Inches(4.8))
    slide2.shapes.add_picture(executed_images[1], Inches(5.1), Inches(1.5), width=Inches(4.8))

    # PAGE 3
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "GOAT Status FINISHED Count"
    slide3.shapes.add_picture(finished_images[0], Inches(0.3), Inches(1.5), width=Inches(4.8))
    slide3.shapes.add_picture(finished_images[1], Inches(5.1), Inches(1.5), width=Inches(4.8))

    # PAGE 4
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = "TDM_GEN_STATUS FINISHED Count"
    slide4.shapes.add_picture(tdm_finished_images[0], Inches(0.3), Inches(1.5), width=Inches(4.8))
    slide4.shapes.add_picture(tdm_finished_images[1], Inches(5.1), Inches(1.5), width=Inches(4.8))

    prs.save("output.pptx")

    print("PPT created successfully with 3 pages")


if __name__ == "__main__":
    main()