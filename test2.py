import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

def create_chart(df, node, column, title, color, filename):

    node_df = df[df['node_id'] == node]

    plt.figure(figsize=(6,4))
    bars = plt.bar(node_df['date'].astype(str), node_df[column], color=color)

    plt.title(title)
    plt.xlabel("Date")
    plt.ylabel("Count")
    plt.xticks(rotation=45)
    plt.grid(axis='y')

    for bar in bars:
        height = bar.get_height()
        plt.text(
            bar.get_x() + bar.get_width()/2,
            height,
            str(height),
            ha='center',
            va='bottom',
            fontweight='bold'
        )

    plt.tight_layout()
    plt.savefig(filename)
    plt.close()

def main():

    df = pd.read_csv("Data.csv")

    df['since'] = pd.to_datetime(df['since'], errors='coerce')
    df['date'] = df['since'].dt.date

    # -----------------------------
    # TOTAL COUNT
    # -----------------------------
    total_df = (
        df.groupby(['date','node_id'])['group_id']
        .count()
        .reset_index()
        .rename(columns={'group_id':'total_count'})
    )

    # -----------------------------
    # EXECUTED
    # -----------------------------
    executed_df = (
        df[df['tasker_status']=="EXECUTED"]
        .groupby(['date','node_id'])['group_id']
        .count()
        .reset_index()
        .rename(columns={'group_id':'executed_count'})
    )

    # -----------------------------
    # GOAT FINISHED
    # -----------------------------
    goat_df = (
        df[df['goat_status']=="FINISHED"]
        .groupby(['date','node_id'])['group_id']
        .count()
        .reset_index()
        .rename(columns={'group_id':'goat_count'})
    )

    # -----------------------------
    # TDM FINISHED
    # -----------------------------
    tdm_df = (
        df[df['tdm_gen_status']=="FINISHED"]
        .groupby(['date','node_id'])['group_id']
        .count()
        .reset_index()
        .rename(columns={'group_id':'tdm_count'})
    )

    nodes = df['node_id'].unique()

    prs = Presentation()

    for node in nodes:

        # Create charts
        create_chart(total_df, node, "total_count",
                     f"{node} Total Tasks", "blue",
                     f"{node}_total.png")

        create_chart(executed_df, node, "executed_count",
                     f"{node} Executed Tasks", "green",
                     f"{node}_executed.png")

        create_chart(goat_df, node, "goat_count",
                     f"{node} GOAT Finished", "orange",
                     f"{node}_goat.png")

        create_chart(tdm_df, node, "tdm_count",
                     f"{node} TDM Finished", "red",
                     f"{node}_tdm.png")

        # -----------------------------
        # CREATE SLIDE
        # -----------------------------
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{node} Day-wise Statistics"

        slide.shapes.add_picture(f"{node}_total.png",
                                 Inches(0.3), Inches(1.5),
                                 width=Inches(4.5))

        slide.shapes.add_picture(f"{node}_executed.png",
                                 Inches(5.2), Inches(1.5),
                                 width=Inches(4.5))

        slide.shapes.add_picture(f"{node}_goat.png",
                                 Inches(0.3), Inches(4.5),
                                 width=Inches(4.5))

        slide.shapes.add_picture(f"{node}_tdm.png",
                                 Inches(5.2), Inches(4.5),
                                 width=Inches(4.5))

    prs.save("Node_Report.pptx")

    print("PPT created successfully!")

if __name__ == "__main__":
    main()