import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches


def create_chart(df, node, filename):

    node_df = df[df['node_id'] == node]

    x = np.arange(len(node_df['date']))
    width = 0.5

    plt.figure(figsize=(10,5))

    bars = plt.bar(
        x,
        node_df['total_free_hours'],
        width,
        color='skyblue',
        label="Free Hours"
    )

    plt.xticks(x, node_df['date'].astype(str), rotation=45)
    plt.xlabel("Date")
    plt.ylabel("Total Free Hours")
    plt.title(f"{node} Free Time per Day")
    plt.legend()
    plt.grid(axis='y')

    # show values on bars
    for bar in bars:
        height = bar.get_height()
        plt.text(
            bar.get_x() + bar.get_width()/2,
            height,
            f"{height:.2f}",
            ha='center',
            va='bottom'
        )

    plt.tight_layout()
    plt.savefig(filename)
    plt.close()


def main():

    df = pd.read_csv("free_time_summary_nodewise.csv")

    df["date"] = pd.to_datetime(df["date"]).dt.date

    nodes = df["node_id"].unique()

    prs = Presentation()

    for node in nodes:

        image_name = f"{node}_free_time.png"

        create_chart(
            df,
            node,
            image_name
        )

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{node} Free Time Summary"

        slide.shapes.add_picture(
            image_name,
            Inches(1),
            Inches(1.5),
            width=Inches(8)
        )

    prs.save("Free_Time_Report.pptx")

    print("PPT created successfully!")


if __name__ == "__main__":
    main()